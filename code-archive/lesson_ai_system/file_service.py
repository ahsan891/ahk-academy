"""
File parsing service — PDF, PPTX, DOCX, TXT.
Returns plain text for sending to the AI pipeline.
"""
import io
from pathlib import Path


def extract_text(content: bytes, filename: str) -> str:
    """Dispatch to the right parser based on file extension."""
    ext = Path(filename).suffix.lower()
    parsers = {
        ".pdf": _pdf,
        ".pptx": _pptx,
        ".ppt": _pptx,
        ".docx": _docx,
        ".txt": _txt,
        ".md": _txt,
    }
    parser = parsers.get(ext, _txt)
    return parser(content)


def extract_pptx_text(content: bytes) -> str:
    return _pptx(content)


# ─── Parsers ──────────────────────────────────────────────────────────────────

def _pdf(content: bytes) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
        return "\n\n".join(parts)
    except ImportError:
        return "[PDF parsing requires pdfplumber — pip install pdfplumber]"
    except Exception as exc:
        return f"[PDF parsing error: {exc}]"


def _pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as exc:
        return f"[PPTX parsing error: {exc}]"


def _docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "[DOCX parsing requires python-docx — pip install python-docx]"
    except Exception as exc:
        return f"[DOCX parsing error: {exc}]"


def _txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")
