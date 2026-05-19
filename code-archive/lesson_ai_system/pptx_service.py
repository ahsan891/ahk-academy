"""
PPTX generation service.
Converts Claude's JSON slide data into a styled PowerPoint file.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Colour palette
C_NAVY   = RGBColor(0x1B, 0x4F, 0x72)   # slide header bar
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # slide background
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)   # body text
C_ACCENT = RGBColor(0x21, 0x8C, 0x74)   # green accent
C_MUTED  = RGBColor(0x7F, 0x8C, 0x8D)   # visual suggestion text


def create_presentation(slide_data: dict, lesson_name: str, output_path: str) -> None:
    """Build and save a .pptx file from structured slide data."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides = slide_data.get("slides", [])
    if not slides:
        raise ValueError("No slides found in slide_data")

    for i, info in enumerate(slides):
        if i == 0:
            _title_slide(prs, lesson_name, info)
        else:
            _content_slide(prs, info)

    prs.save(output_path)


# ─── Slide builders ───────────────────────────────────────────────────────────

def _title_slide(prs: Presentation, lesson_name: str, info: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, C_NAVY)

    # Main title
    _textbox(
        slide,
        text=info.get("title", lesson_name),
        left=Inches(0.8), top=Inches(2.0),
        width=Inches(11.7), height=Inches(1.6),
        size=Pt(44), bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Subtitle (first bullet)
    bullets = info.get("content", [])
    if bullets:
        _textbox(
            slide,
            text=bullets[0],
            left=Inches(1.5), top=Inches(3.8),
            width=Inches(10.3), height=Inches(1.0),
            size=Pt(22), bold=False,
            color=RGBColor(0xA9, 0xCC, 0xE3),
            align=PP_ALIGN.CENTER,
        )

    # Decorative accent bar
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.07),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    _add_notes(slide, info.get("speaker_notes", ""))


def _content_slide(prs: Presentation, info: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_LIGHT)

    # Header bar
    hbar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.25),
    )
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = C_NAVY
    hbar.line.fill.background()

    # Title text in header
    _textbox(
        slide,
        text=info.get("title", ""),
        left=Inches(0.35), top=Inches(0.12),
        width=Inches(12.6), height=Inches(1.0),
        size=Pt(26), bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Bullet points
    bullets = info.get("content", [])
    if bullets:
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.6),
        )
        tf = box.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(10)
            p.space_after  = Pt(4)
            run = p.add_run()
            run.text = f"▸  {bullet}"
            run.font.size  = Pt(20)
            run.font.color.rgb = C_TEXT

    # Visual suggestion (bottom-right corner)
    visual = info.get("visual_suggestion", "")
    if visual:
        _textbox(
            slide,
            text=f"🖼  {visual[:120]}",
            left=Inches(7.8), top=Inches(6.1),
            width=Inches(5.2), height=Inches(1.1),
            size=Pt(11), bold=False, color=C_MUTED,
            align=PP_ALIGN.LEFT, italic=True,
        )

    _add_notes(slide, info.get("speaker_notes", ""))


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(
    slide, text: str,
    left, top, width, height,
    size: Pt, bold: bool = False,
    color: RGBColor = C_TEXT,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color


def _add_notes(slide, notes: str) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
